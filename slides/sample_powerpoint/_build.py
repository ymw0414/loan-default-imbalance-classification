"""
Sample 1: Academic-style PowerPoint
Clean, classic look. Navy + accent blue. Designed for 10-min talk
with two presenters (Minwoo 5 min + Nathan 5 min).
"""
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE         = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(HERE, '..', '..'))
CV_DIR       = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'cv_results', 'GridSearchCV')
COMBINED_CSV = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'combined_model_comparison.csv')
IMBAL_CSV    = os.path.join(CV_DIR, 'imbalance_compare.csv')
FIG_DIR      = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'figure')

OUT = os.path.join(HERE, 'sample_academic.pptx')

# --------------------------- color palette (academic) ---------------------------
NAVY    = RGBColor(0x1F, 0x3A, 0x68)
ACCENT  = RGBColor(0x2E, 0x86, 0xC1)
GREY    = RGBColor(0x4A, 0x4A, 0x4A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xEC, 0xF0, 0xF1)
RED     = RGBColor(0xC0, 0x39, 0x2B)


def fmt(v):
    return 'TBD' if pd.isna(v) else f'{v:.4f}'


def set_text(tf, text, *, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color


def add_p(tf, text, *, size=14, bold=False, color=None, level=0, bullet=True):
    p = tf.add_paragraph()
    p.text = ('• ' if bullet else '') + text
    p.level = level
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color


def title(slide, text, speaker=''):
    """Section title in top-left, speaker tag top-right."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.5), Inches(0.7))
    set_text(tb.text_frame, text, size=26, bold=True, color=NAVY)
    if speaker:
        tb = slide.shapes.add_textbox(Inches(10.0), Inches(0.4), Inches(3.0), Inches(0.5))
        tf = tb.text_frame
        tf.text = speaker
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(11)
            r.font.italic = True
            r.font.color.rgb = ACCENT


def add_table(slide, df, top, left, width, height, col_widths=None, header_color=NAVY):
    rows, cols = df.shape[0] + 1, df.shape[1]
    t = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    if col_widths:
        for j, w in enumerate(col_widths):
            t.columns[j].width = Inches(w)
    for j, name in enumerate(df.columns):
        cell = t.cell(0, j)
        cell.text = str(name)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = WHITE
    for i in range(rows - 1):
        for j in range(cols):
            cell = t.cell(i + 1, j)
            v = df.iloc[i, j]
            cell.text = '' if pd.isna(v) else str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
    return t


# --------------------------- load results ---------------------------
combined  = pd.read_csv(COMBINED_CSV) if os.path.exists(COMBINED_CSV) else None
imbalance = pd.read_csv(IMBAL_CSV) if os.path.exists(IMBAL_CSV) else None

best_label = combined['model'].iloc[0] if combined is not None else 'TBD'
best_auc   = fmt(combined['best_val_AUC'].iloc[0]) if combined is not None else 'TBD'

# --------------------------- build deck ---------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# --- Slide 1: Title (Minwoo, 30s) ---
s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(1.4))
set_text(tb.text_frame, 'Loan Default Prediction', size=44, bold=True, color=NAVY)

tb = s.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(1))
set_text(tb.text_frame, 'Comparing Classic ML and Neural Networks under Class Imbalance',
         size=22, color=GREY)

tb = s.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(12), Inches(0.4))
set_text(tb.text_frame, 'Minwoo Yoo  ·  Nathan',
         size=16, color=GREY)
tb = s.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(12), Inches(0.4))
set_text(tb.text_frame, 'Home Credit Default Risk · Kaggle',
         size=13, color=GREY)


# --- Slide 2: Problem & Motivation (Minwoo, 1 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Problem & Motivation', '📍 Minwoo · 1 min')

tb = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.13), Inches(5.6))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'The task', size=18, bold=True, color=NAVY)
add_p(tf, 'Predict the probability that a loan applicant defaults  (binary, TARGET ∈ {0, 1}).', size=14)
add_p(tf, '', bullet=False)
add_p(tf, 'Why class imbalance matters', size=18, bold=True, color=NAVY, bullet=False)
add_p(tf, 'Only ~8.07% of applicants default — the data is severely imbalanced.', size=14)
add_p(tf, 'A "predict-no-default" baseline scores 92% accuracy with zero recall.', size=14)
add_p(tf, 'A model can hit ROC AUC = 0.745 yet classify zero defaulters at the 0.5 threshold.', size=14)
add_p(tf, '', bullet=False)
add_p(tf, 'Research question', size=18, bold=True, color=NAVY, bullet=False)
add_p(tf, 'Do classic shallow learners and neural networks respond differently to class imbalance handling?', size=15, color=ACCENT)


# --- Slide 3: Dataset & EDA (Minwoo, 1 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Dataset · Home Credit Default Risk', '📍 Minwoo · 1 min')

# Left: text
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.7), Inches(5.8))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'Size', size=17, bold=True, color=NAVY)
add_p(tf, 'Train: 307,511 applicants × 122 features.', size=13)
add_p(tf, 'Test: 48,744 applicants (Kaggle hidden labels).', size=13)
add_p(tf, '', bullet=False)
add_p(tf, 'Class imbalance', size=17, bold=True, color=NAVY, bullet=False)
add_p(tf, 'Default rate ≈ 8.07% (24,825 / 307,511).', size=13)
add_p(tf, '', bullet=False)
add_p(tf, 'Strongest signals', size=17, bold=True, color=NAVY, bullet=False)
add_p(tf, 'EXT_SOURCE_1/2/3 (external credit scores), |ρ| ≈ 0.16–0.18.', size=13)
add_p(tf, 'DAYS_EMPLOYED sentinel 365,243 flags 18% (unemployed).', size=13)
add_p(tf, 'Building-info columns are 50–70% missing.', size=13)

# Right: class balance chart
chart_path = os.path.join(FIG_DIR, 'class_balance.png')
if os.path.exists(chart_path):
    s.shapes.add_picture(chart_path, Inches(7.3), Inches(2.0), width=Inches(5.8))


# --- Slide 4: Methodology / Pipeline (Nathan, 1 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Methodology · Shared Pipeline', '📍 Nathan · 1 min')

tb = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.13), Inches(5.8))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'Preprocessing (identical for all five models)', size=18, bold=True, color=NAVY)
add_p(tf, '80% / 20% stratified train–val split (random_state = 42)', size=13)
add_p(tf, 'Drop SK_ID_CURR identifier  ·  keep IDs only for the submission file', size=13)
add_p(tf, 'Mean imputation for numeric NaN  ·  one-hot encode 16 categoricals → 244 features', size=13)
add_p(tf, 'StandardScaler fitted on train only (essential for LR / MLP)', size=13)
add_p(tf, '', bullet=False)
set_text(tf := s.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12.13), Inches(2.5)).text_frame,
         'Tuning & validation', size=18, bold=True, color=NAVY)
add_p(tf, 'GridSearchCV with PredefinedSplit (single train→val fold, no random k-fold)', size=13)
add_p(tf, 'Scoring: roc_auc  ·  best parameter setting picked per model', size=13)
add_p(tf, 'CV CSVs written to result/home_credit/cv_results/GridSearchCV/', size=13)


# --- Slide 5: Five Models Compared (Nathan, 1.5 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Five Models Compared', '📍 Nathan · 1.5 min')

df_models = pd.DataFrame([
    ['Logistic Regression',           'Linear',          'Nathan',  'C, tol  ·  class_weight = balanced'],
    ['Random Forest',                  'Tree (bagging)',  'Nathan',  'n_estimators, min_samples_*  ·  class_weight'],
    ['HistGradientBoosting',           'Tree (boosting)', 'Nathan',  'learning_rate, max_iter, min_samples_leaf'],
    ['Shallow MLP — hidden (64,)',     'Neural network',  'Minwoo',  'alpha, learning_rate_init  ·  early stopping'],
    ['Deep MLP — hidden (128, 64, 32)','Neural network',  'Minwoo',  'alpha, learning_rate_init  ·  early stopping'],
], columns=['Model', 'Family', 'Owner', 'Tuned hyperparameters'])
add_table(s, df_models, top=1.3, left=0.5, width=12.3, height=2.8,
          col_widths=[3.4, 2.2, 1.4, 5.3])

tb = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'Why this lineup?', size=17, bold=True, color=NAVY)
add_p(tf, 'Spans the course curriculum: linear → tree (bagging + boosting) → shallow + deep neural networks.', size=13)
add_p(tf, 'Lets us isolate whether NN behaves differently from classic shallow learners under imbalance.', size=13)
add_p(tf, '', bullet=False)
add_p(tf, 'Same preprocessing, same train/val split, same scoring → fair comparison.', size=13, bold=True, color=ACCENT)


# --- Slide 6: Combined Validation AUC (Nathan, 1 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Combined Validation ROC AUC', '📍 Nathan · 1 min')

chart_path = os.path.join(FIG_DIR, 'model_auc_bar.png')
if os.path.exists(chart_path):
    s.shapes.add_picture(chart_path, Inches(0.4), Inches(1.2), width=Inches(7.6))

if combined is not None:
    cmb = combined.copy()
    cmb['best_val_AUC'] = cmb['best_val_AUC'].apply(fmt)
    cmb = cmb[['model', 'best_val_AUC']]
    cmb.columns = ['Model', 'AUC']
    cmb['Model'] = cmb['Model'].astype(str).str.replace(r'\s*\(.*\)$', '', regex=True)
    add_table(s, cmb, top=1.4, left=8.3, width=4.7, height=3.0, col_widths=[3.0, 1.7])

tb = s.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, f'Best model: {best_label}  (val ROC AUC = {best_auc})',
         size=17, bold=True, color=NAVY)
add_p(tf, 'Tree boosting wins by ~1 AUC point — consistent with the literature on tabular credit data.', size=13)
add_p(tf, 'Neural networks tie LR and Random Forest; depth alone does not unlock new performance.', size=13)


# --- Slide 7: Neural Network Track (Minwoo, 1.5 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Neural Network Track', '📍 Minwoo · 1.5 min')

tb = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(7.0), Inches(5.8))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'Architecture', size=17, bold=True, color=NAVY)
add_p(tf, 'Shallow MLP — hidden (64,)  ·  ~15.7K parameters', size=13)
add_p(tf, 'Deep MLP — hidden (128, 64, 32)  ·  ~42K parameters', size=13)
add_p(tf, 'ReLU activation, Adam optimizer, early_stopping = True', size=13)
add_p(tf, '', bullet=False)
set_text(tf, 'Hyperparameter tuning', size=17, bold=True, color=NAVY)
# Workaround: using add_p again to continue
add_p(tf, '', bullet=False)
add_p(tf, 'alpha (L2 reg.) ∈ {1e-4, 1e-3}  ·  learning_rate_init ∈ {0.001, 0.005}', size=13)
add_p(tf, 'Best for both: alpha = 1e-3, lr = 0.005', size=13)
add_p(tf, '', bullet=False)
set_text(tf, 'Findings', size=17, bold=True, color=NAVY)
add_p(tf, 'Shallow 0.7405  vs  Deep 0.7416  → depth adds ~0.001 AUC (noise).', size=13)
add_p(tf, 'NN ranks 4th/5th overall; bottleneck is the feature set, not the model.', size=13)

# Right: place a small custom comparison
df_nn = pd.DataFrame([
    ['Shallow MLP', '(64,)',         '0.7405'],
    ['Deep MLP',    '(128, 64, 32)', '0.7416'],
], columns=['Model', 'Hidden layers', 'Val AUC'])
add_table(s, df_nn, top=1.4, left=7.8, width=5.2, height=1.4, col_widths=[1.6, 1.8, 1.2])


# --- Slide 8: Imbalance Handling Deep Dive (Minwoo, 1.5 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Imbalance Handling on the Best MLP', '📍 Minwoo · 1.5 min')

if imbalance is not None and len(imbalance) > 0:
    tbl_df = imbalance[['method', 'threshold', 'ROC_AUC', 'PR_AUC', 'precision', 'recall', 'F1']].copy()
    add_table(s, tbl_df, top=1.2, left=0.5, width=7.5, height=1.7,
              col_widths=[1.2, 1.0, 1.1, 1.0, 1.1, 1.0, 1.1])

    chart_path = os.path.join(FIG_DIR, 'imbalance_compare.png')
    if os.path.exists(chart_path):
        s.shapes.add_picture(chart_path, Inches(0.4), Inches(3.0), width=Inches(7.4))

# Right side commentary
tb = s.shapes.add_textbox(Inches(8.2), Inches(1.2), Inches(4.9), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'Three strategies', size=17, bold=True, color=NAVY)
add_p(tf, 'Baseline — train as is, classify at 0.5', size=12)
add_p(tf, 'SMOTE — synthetic minority oversampling', size=12)
add_p(tf, 'Threshold tuning — pick F1-optimal cutoff', size=12)
add_p(tf, '', bullet=False)
set_text(tf, 'Headline finding', size=17, bold=True, color=NAVY)
add_p(tf, '', bullet=False)
add_p(tf, 'Baseline: AUC 0.745, recall 0% (no prediction crosses 0.5)', size=12)
add_p(tf, 'SMOTE hurt: AUC 0.745 → 0.648', size=12, color=RED)
add_p(tf, 'Threshold tuning at 0.164: recall 40%, F1 0.30 ✓', size=12, color=ACCENT, bold=True)
add_p(tf, '', bullet=False)
add_p(tf, 'Choosing the threshold matters more than rebalancing the data.', size=13, bold=True, color=NAVY)


# --- Slide 9: Conclusions (Both, 1 min) ---
s = prs.slides.add_slide(blank)
title(s, 'Conclusions', '📍 Both · 1 min')

tb = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.13), Inches(5.6))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'What we learned', size=18, bold=True, color=NAVY)
add_p(tf, f'Best model: {best_label} at val ROC AUC = {best_auc}; tree boosting wins by ~1 AUC point.', size=14)
add_p(tf, 'Neural networks tie LR and RF; going deeper does not help on this tabular data.', size=14)
add_p(tf, 'On the MLP, threshold tuning lifts recall from 0% to 40% — SMOTE hurt AUC.', size=14, bold=True, color=ACCENT)
add_p(tf, '', bullet=False)
set_text(tf, 'Limitations', size=18, bold=True, color=NAVY)
add_p(tf, '', bullet=False)
add_p(tf, 'Only main application_*.csv tables — auxiliary tables typically add several AUC points.', size=13)
add_p(tf, 'PredefinedSplit is fast but optimistic; full k-fold CV would tighten confidence intervals.', size=13)
add_p(tf, '', bullet=False)
set_text(tf, 'Future work', size=18, bold=True, color=NAVY)
add_p(tf, '', bullet=False)
add_p(tf, 'Integrate auxiliary tables  ·  ensemble HGBC + best MLP  ·  Platt/isotonic calibration', size=13)


# --- Slide 10: Q&A / Closing (Both, 30s) ---
s = prs.slides.add_slide(blank)
title(s, 'Thank You', '📍 Both · 30s')

tb = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.0))
set_text(tb.text_frame, 'Questions?', size=42, bold=True, color=NAVY)

tb = s.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, 'Recording link', size=18, bold=True, color=ACCENT)
add_p(tf, '<<INSERT RECORDING LINK BEFORE SUBMISSION>>', size=14)
add_p(tf, '', bullet=False)
set_text(tf, 'Code', size=18, bold=True, color=ACCENT)
add_p(tf, '', bullet=False)
add_p(tf, 'github.com/ymw0414/loan-default-imbalance-classification', size=14)


# --------------------------- save ---------------------------
prs.save(OUT)
print(f'wrote {OUT}  ·  {len(prs.slides)} slides')
