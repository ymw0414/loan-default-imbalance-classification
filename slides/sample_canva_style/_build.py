"""
Sample 3: Canva-style modern PowerPoint

Visual recipe (mimics Canva's modern look):
  - Large colored side panels / accent shapes
  - Bold display typography for slide titles (oversize)
  - Plenty of whitespace
  - Minimal bullets, big numbers as visual anchors
  - Limited but high-contrast color palette

Result: PPTX that feels like a Canva template but is actually a
plain .pptx (no Canva account / web access required).
"""
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

HERE         = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(HERE, '..', '..'))
CV_DIR       = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'cv_results', 'GridSearchCV')
COMBINED_CSV = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'combined_model_comparison.csv')
IMBAL_CSV    = os.path.join(CV_DIR, 'imbalance_compare.csv')
FIG_DIR      = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'figure')

OUT = os.path.join(HERE, 'sample_canva_style.pptx')

# --------------------------- Canva-ish palette ---------------------------
# Bolder, more saturated. Inspired by modern Canva tabular templates.
DEEP    = RGBColor(0x14, 0x29, 0x4B)   # very dark navy (background block)
TEAL    = RGBColor(0x1A, 0xBC, 0x9C)   # teal accent
ORANGE  = RGBColor(0xF3, 0x9C, 0x12)   # warm accent
RED     = RGBColor(0xE7, 0x4C, 0x3C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_W  = RGBColor(0xFA, 0xFA, 0xFA)
SOFT    = RGBColor(0xF1, 0xF6, 0xFA)
GREY    = RGBColor(0x55, 0x65, 0x70)
DARK    = RGBColor(0x2C, 0x3E, 0x50)


def fmt(v):
    return 'TBD' if pd.isna(v) else f'{v:.4f}'


def fill_bg(slide, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide.part.package.presentation_part.presentation.slide_width, slide.part.package.presentation_part.presentation.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_rect(slide, left, top, w, h, color, no_line=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    if no_line:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tf


def add_bullets(slide, x, y, w, h, lines, *, size=14, color=DARK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        if i == 0:
            tf.text = ln
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = ln
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tf


def speaker_chip(slide, x, y, text, color=TEAL):
    """A small rounded color chip that says e.g. 'Minwoo · 1 min'."""
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(0.4))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    chip.text_frame.text = text
    p = chip.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE


def big_number(slide, x, y, w, h, number, label, color=TEAL):
    """Big metric: number on top, small label below."""
    add_text(slide, x, y, w, 1.3, number, size=54, bold=True, color=color)
    add_text(slide, x, y + 1.2, w, 0.4, label, size=12, color=GREY)


# --------------------------- load results ---------------------------
combined  = pd.read_csv(COMBINED_CSV) if os.path.exists(COMBINED_CSV) else None
imbalance = pd.read_csv(IMBAL_CSV) if os.path.exists(IMBAL_CSV) else None

best_label = combined['model'].iloc[0] if combined is not None else 'TBD'
best_auc   = fmt(combined['best_val_AUC'].iloc[0]) if combined is not None else 'TBD'

# --------------------------- build deck ---------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

W, H = 13.333, 7.5


# --- Slide 1: Title ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, DEEP)              # full deep-navy bg
add_rect(s, 0, 0, 0.4, H, TEAL)            # left teal stripe
# big title
add_text(s, 0.9, 2.0, W - 1.8, 1.2,
         'Loan Default', size=64, bold=True, color=WHITE)
add_text(s, 0.9, 3.0, W - 1.8, 1.0,
         'Prediction', size=64, bold=True, color=TEAL)
add_text(s, 0.9, 4.5, W - 1.8, 0.5,
         'Comparing Classic ML and Neural Networks under Class Imbalance',
         size=18, color=NEAR_W)
add_text(s, 0.9, 6.5, W - 1.8, 0.4,
         'Minwoo Yoo · Nathan      |      Home Credit Default Risk · Kaggle',
         size=12, italic=True, color=NEAR_W)


# --- Slide 2: Problem & Motivation ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, W, 1.2, DEEP)            # dark header
add_text(s, 0.6, 0.32, W - 4, 0.7, 'Why this problem?', size=28, bold=True, color=WHITE)
speaker_chip(s, W - 3.2, 0.42, 'Minwoo · 1 min', TEAL)

# The big number anchor
big_number(s, 0.6, 1.7, 4, 2, '8.07%', 'of applicants actually default', TEAL)
big_number(s, 0.6, 4.3, 4, 2, '0%', 'recall at 0.5 threshold (baseline MLP)', RED)

# Bullet body on the right
add_text(s, 5.2, 1.7, 7.5, 0.5,
         'Lenders need to rank applicants by default risk', size=18, bold=True, color=DARK)
add_bullets(s, 5.2, 2.4, 7.5, 4,
            ['• Imbalance: 92% non-default vs 8% default',
             '• "Always predict 0" → 92% accuracy, 0 recall — useless',
             '• Even a model with ROC AUC = 0.745 can classify 0 defaulters',
             '  if every prediction sits below 0.5',
             '',
             'Research question:',
             ''],
            size=14)
add_text(s, 5.2, 5.9, 7.5, 0.8,
         'Do classic shallow learners and NNs respond differently to imbalance handling?',
         size=15, italic=True, color=ORANGE, bold=True)


# --- Slide 3: Dataset & EDA ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, W, 1.2, DEEP)
add_text(s, 0.6, 0.32, W - 4, 0.7, 'Dataset · Home Credit', size=28, bold=True, color=WHITE)
speaker_chip(s, W - 3.2, 0.42, 'Minwoo · 1 min', TEAL)

# Three big number callouts
big_number(s, 0.6, 1.7, 3.5, 2, '307K', 'applicants in train', TEAL)
big_number(s, 4.4, 1.7, 3.5, 2, '244', 'features after one-hot', ORANGE)
big_number(s, 8.2, 1.7, 4.5, 2, '8.07%', 'positive class', RED)

# Image at the bottom
chart_path = os.path.join(FIG_DIR, 'class_balance.png')
if os.path.exists(chart_path):
    s.shapes.add_picture(chart_path, Inches(0.6), Inches(4.2), width=Inches(7.0))

add_text(s, 8.0, 4.2, 5.0, 0.5, 'Strongest signals', size=18, bold=True, color=DARK)
add_bullets(s, 8.0, 4.8, 5.0, 2.5,
            ['EXT_SOURCE_1/2/3 (external credit scores)',
             'DAYS_EMPLOYED 365,243 = unemployed flag',
             'Building-info columns 50–70% missing'], size=13)


# --- Slide 4: Methodology ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, W, 1.2, DEEP)
add_text(s, 0.6, 0.32, W - 4, 0.7, 'Shared pipeline', size=28, bold=True, color=WHITE)
speaker_chip(s, W - 3.2, 0.42, 'Nathan · 1 min', ORANGE)

# Pipeline steps as connected colored cards
steps = [
    ('Split', '80% / 20%\nstratified', TEAL),
    ('Drop ID', 'remove\nSK_ID_CURR', TEAL),
    ('Impute', 'mean for\nnumeric NaN', TEAL),
    ('Encode', 'one-hot →\n244 features', ORANGE),
    ('Scale', 'StandardScaler\non train only', ORANGE),
    ('Tune', 'GridSearchCV\nROC AUC', RED),
]
card_w = 1.85
gap = 0.15
start_x = (W - (len(steps)*card_w + (len(steps)-1)*gap)) / 2
for i, (head, body, c) in enumerate(steps):
    x = start_x + i * (card_w + gap)
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(card_w), Inches(2.5))
    rect.fill.solid(); rect.fill.fore_color.rgb = c
    rect.line.fill.background()
    rect.shadow.inherit = False
    add_text(s, x, 2.5, card_w, 0.6, head, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, 3.3, card_w, 1.4, body, size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 5.6, W - 1.2, 0.5, 'Same preprocessing for all five models — fair comparison.',
         size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 6.2, W - 1.2, 0.5,
         'PredefinedSplit (single train→val fold) · scoring = roc_auc · seed = 42',
         size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# --- Slide 5: Five Models Compared ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, W, 1.2, DEEP)
add_text(s, 0.6, 0.32, W - 4, 0.7, 'Five models', size=28, bold=True, color=WHITE)
speaker_chip(s, W - 3.2, 0.42, 'Nathan · 1.5 min', ORANGE)

models = [
    ('Logistic Regression',           'Linear',           'Nathan',  TEAL),
    ('Random Forest',                  'Tree (bagging)',   'Nathan',  TEAL),
    ('HistGradBoost',                  'Tree (boosting)',  'Nathan',  TEAL),
    ('Shallow MLP (64,)',              'Neural network',   'Minwoo',  ORANGE),
    ('Deep MLP (128, 64, 32)',         'Neural network',   'Minwoo',  ORANGE),
]
card_w = 2.4
gap = 0.15
start_x = (W - (len(models)*card_w + (len(models)-1)*gap)) / 2
for i, (name, fam, owner, c) in enumerate(models):
    x = start_x + i * (card_w + gap)
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(card_w), Inches(3.4))
    rect.fill.solid(); rect.fill.fore_color.rgb = c
    rect.line.fill.background()
    rect.shadow.inherit = False
    add_text(s, x, 2.3, card_w, 0.8, name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, 3.4, card_w, 0.5, fam, size=11, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, x, 4.3, card_w, 0.5, f'by {owner}', size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 6.0, W - 1.2, 0.5,
         'Spans the curriculum: linear → tree (bagging + boosting) → shallow + deep NN.',
         size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# --- Slide 6: Combined Validation AUC ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, W, 1.2, DEEP)
add_text(s, 0.6, 0.32, W - 4, 0.7, 'Combined ROC AUC', size=28, bold=True, color=WHITE)
speaker_chip(s, W - 3.2, 0.42, 'Nathan · 1 min', ORANGE)

chart_path = os.path.join(FIG_DIR, 'model_auc_bar.png')
if os.path.exists(chart_path):
    s.shapes.add_picture(chart_path, Inches(0.4), Inches(1.5), width=Inches(7.6))

# Right side big number callouts
big_number(s, 8.5, 1.6, 4.5, 2.0, best_auc, f'{best_label}', TEAL)
add_text(s, 8.5, 3.8, 4.5, 0.5, 'Headline', size=16, bold=True, color=DARK)
add_bullets(s, 8.5, 4.4, 4.5, 2.5,
            ['Tree boosting wins by ~1 AUC point',
             'Consistent with literature on tabular',
             'NN ties LR & RF — not the model, the features'],
            size=13)


# --- Slide 7: Neural Network Track ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, W, 1.2, DEEP)
add_text(s, 0.6, 0.32, W - 4, 0.7, 'Neural networks', size=28, bold=True, color=WHITE)
speaker_chip(s, W - 3.2, 0.42, 'Minwoo · 1.5 min', TEAL)

# Two big architecture cards
add_rect(s, 0.6, 1.7, 5.8, 4.0, ORANGE)
add_text(s, 0.8, 1.9, 5.4, 0.6, 'Shallow MLP', size=22, bold=True, color=WHITE)
add_text(s, 0.8, 2.6, 5.4, 0.5, '(64,)  ·  ~15.7K params', size=14, color=WHITE, italic=True)
big_number(s, 0.8, 3.4, 5.4, 1.5, '0.7405', 'val ROC AUC', WHITE)

add_rect(s, 6.9, 1.7, 5.8, 4.0, RED)
add_text(s, 7.1, 1.9, 5.4, 0.6, 'Deep MLP', size=22, bold=True, color=WHITE)
add_text(s, 7.1, 2.6, 5.4, 0.5, '(128, 64, 32)  ·  ~42K params', size=14, color=WHITE, italic=True)
big_number(s, 7.1, 3.4, 5.4, 1.5, '0.7416', 'val ROC AUC', WHITE)

add_text(s, 0.6, 6.0, W - 1.2, 0.5,
         '+0.001 with 3× the parameters → depth alone does not help on this tabular data.',
         size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 6.6, W - 1.2, 0.5,
         'Best for both: alpha = 1e-3, learning_rate_init = 0.005',
         size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# --- Slide 8: Imbalance Handling ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, W, 1.2, DEEP)
add_text(s, 0.6, 0.32, W - 4, 0.7, 'Imbalance handling', size=28, bold=True, color=WHITE)
speaker_chip(s, W - 3.2, 0.42, 'Minwoo · 1.5 min', TEAL)

# Three result cards
add_rect(s, 0.6, 1.6, 4.0, 3.7, GREY)
add_text(s, 0.6, 1.8, 4.0, 0.5, 'Baseline', size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
big_number(s, 0.6, 2.4, 4.0, 1.5, '0%', 'recall @ 0.5', WHITE)
add_text(s, 0.6, 4.5, 4.0, 0.5, 'AUC 0.745', size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 4.8, 1.6, 4.0, 3.7, RED)
add_text(s, 4.8, 1.8, 4.0, 0.5, 'SMOTE ✗', size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
big_number(s, 4.8, 2.4, 4.0, 1.5, '0.648', 'AUC dropped', WHITE)
add_text(s, 4.8, 4.5, 4.0, 0.5, 'recall 21%, F1 0.18', size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 9.0, 1.6, 4.0, 3.7, TEAL)
add_text(s, 9.0, 1.8, 4.0, 0.5, 'Threshold ✓', size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
big_number(s, 9.0, 2.4, 4.0, 1.5, '40%', 'recall @ 0.164', WHITE)
add_text(s, 9.0, 4.5, 4.0, 0.5, 'AUC 0.745, F1 0.30', size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 5.7, W - 1.2, 0.7,
         'Choosing the threshold matters more than rebalancing the data.',
         size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 6.5, W - 1.2, 0.4,
         'Same model, different decision rule — recall jumps from 0% to 40%.',
         size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# --- Slide 9: Conclusions ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, DEEP)
add_text(s, 0.7, 0.5, W - 1.4, 1.0, 'Takeaways', size=44, bold=True, color=WHITE)
add_rect(s, 0.7, 1.6, 1.0, 0.06, TEAL)

# 3 boxes
boxes = [
    (1, 'Tree boosting wins',
        f'{best_label} at {best_auc} val AUC.\nTabular data favors GBT over NN.', TEAL),
    (2, 'Threshold > resampling',
        'On the MLP, threshold tuning lifted recall\nfrom 0% to 40% — SMOTE hurt AUC.', ORANGE),
    (3, 'Depth ≈ noise',
        'Shallow vs Deep MLP differ by 0.001 AUC.\nFeatures are the bottleneck.', RED),
]
for i, (num, head, body, c) in enumerate(boxes):
    x = 0.7 + i * 4.2
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(4.0), Inches(3.6))
    rect.fill.solid(); rect.fill.fore_color.rgb = c
    rect.line.fill.background()
    rect.shadow.inherit = False
    add_text(s, x, 2.4, 4.0, 0.6, f'#{num}', size=22, bold=True, color=WHITE)
    add_text(s, x, 3.0, 4.0, 0.6, head, size=18, bold=True, color=WHITE)
    add_text(s, x, 3.7, 4.0, 1.6, body, size=13, color=WHITE)

add_text(s, 0.7, 6.2, W - 1.4, 0.5,
         'Limitations: only main tables (auxiliary tables → +2–3 AUC pts) · single split CV.',
         size=12, italic=True, color=NEAR_W)


# --- Slide 10: Q&A ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, NEAR_W)
add_rect(s, 0, 0, 4.5, H, DEEP)
add_text(s, 0.5, 2.7, 4.0, 1.5, 'Q&A', size=88, bold=True, color=WHITE)
add_text(s, 0.5, 4.3, 4.0, 0.5, 'thank you', size=18, italic=True, color=TEAL)

# Right side info
add_text(s, 5.2, 2.5, 7.5, 0.5, 'Recording link', size=18, bold=True, color=ORANGE)
add_text(s, 5.2, 3.0, 7.5, 0.5,
         '<<INSERT RECORDING LINK BEFORE SUBMISSION>>',
         size=14, color=DARK)

add_text(s, 5.2, 4.0, 7.5, 0.5, 'Code', size=18, bold=True, color=ORANGE)
add_text(s, 5.2, 4.5, 7.5, 0.5,
         'github.com/ymw0414/loan-default-imbalance-classification',
         size=14, color=DARK)

add_text(s, 5.2, 5.6, 7.5, 0.5, 'Team', size=18, bold=True, color=ORANGE)
add_text(s, 5.2, 6.1, 7.5, 0.5,
         'Minwoo Yoo · Nathan',
         size=14, color=DARK)


# --------------------------- save ---------------------------
prs.save(OUT)
print(f'wrote {OUT}  ·  {len(prs.slides)} slides')
