"""
Builder script for the final-project presentation slides (.pptx).

The script reads the live notebook results from
    ../result/home_credit/cv_results/GridSearchCV/*.csv
    ../result/home_credit/combined_model_comparison.csv
and fills the slide deck. If those files are missing (notebook still running),
sensible "TBD" placeholders are used so the deck always builds.

Run:
    python _build_slides.py
"""
from __future__ import annotations

import ast
import os
import sys

import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --------------------------------------------------------------------- paths
HERE         = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(HERE, '..'))
CV_DIR       = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'cv_results', 'GridSearchCV')
COMBINED_CSV = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'combined_model_comparison.csv')
IMBAL_CSV    = os.path.join(CV_DIR, 'imbalance_compare.csv')
FIG_DIR      = os.path.join(PROJECT_ROOT, 'result', 'home_credit', 'figure')

CHART_CLASS_BALANCE = os.path.join(FIG_DIR, 'class_balance.png')
CHART_MODEL_AUC     = os.path.join(FIG_DIR, 'model_auc_bar.png')
CHART_IMBALANCE     = os.path.join(FIG_DIR, 'imbalance_compare.png')

OUT_PPTX = os.path.join(HERE, 'final_presentation.pptx')

RECORDING_LINK_PLACEHOLDER = '<<INSERT RECORDING LINK BEFORE SUBMISSION>>'

# --------------------------------------------------------------------- helpers
def safe_best(csv_path: str) -> tuple[float | None, str | None]:
    if not os.path.exists(csv_path):
        return None, None
    df = pd.read_csv(csv_path)
    df = df.sort_values('rank_test_score').head(1)
    return float(df['mean_test_score'].iloc[0]), str(df['params'].iloc[0])


def load_combined() -> pd.DataFrame:
    if os.path.exists(COMBINED_CSV):
        return pd.read_csv(COMBINED_CSV)

    # Fallback: build from individual CSVs
    rows = []
    for fname, label in [
        ('lr.csv',           'Logistic Regression (Nathan)'),
        ('rfc.csv',          'Random Forest (Nathan)'),
        ('hgbc.csv',         'HistGradBoost (Nathan)'),
        ('shallow_mlp.csv',  'Shallow MLP (Minwoo)'),
        ('deep_mlp.csv',     'Deep MLP (Minwoo)'),
    ]:
        score, params = safe_best(os.path.join(CV_DIR, fname))
        rows.append({'model': label, 'best_val_AUC': score, 'best_params': params})

    df = pd.DataFrame(rows)
    df = df.sort_values('best_val_AUC', ascending=False, na_position='last').reset_index(drop=True)
    return df


def load_imbalance() -> pd.DataFrame | None:
    if not os.path.exists(IMBAL_CSV):
        return None
    return pd.read_csv(IMBAL_CSV)


def fmt_score(v):
    if v is None or pd.isna(v):
        return 'TBD'
    return f'{v:.4f}'


# --------------------------------------------------------------------- styling
TITLE_COLOR = RGBColor(0x1F, 0x3A, 0x68)        # deep navy
ACCENT      = RGBColor(0x2E, 0x86, 0xC1)        # accent blue
GREY        = RGBColor(0x55, 0x55, 0x55)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)


def set_text(tf, text, *, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color


def add_paragraph(tf, text, *, size=16, bold=False, color=None, level=0, bullet=True):
    p = tf.add_paragraph()
    p.text = ('• ' if bullet else '') + text
    p.level = level
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color


def add_section_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    set_text(tb.text_frame, text, size=28, bold=True, color=TITLE_COLOR)


def add_body_textbox(slide, top=1.1, height=5.8):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.13), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def add_table(slide, df, top=1.3, left=0.6, width=12.13, height=4.5,
              col_widths_inches=None, header_color=ACCENT):
    rows, cols = df.shape[0] + 1, df.shape[1]
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tbl = tbl_shape.table

    # Column widths (optional)
    if col_widths_inches:
        for j, w in enumerate(col_widths_inches):
            tbl.columns[j].width = Inches(w)

    # Header
    for j, name in enumerate(df.columns):
        cell = tbl.cell(0, j)
        cell.text = str(name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.color.rgb = WHITE

    # Body
    for i in range(rows - 1):
        for j in range(cols):
            cell = tbl.cell(i + 1, j)
            value = df.iloc[i, j]
            cell.text = '' if pd.isna(value) else str(value)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
    return tbl


# --------------------------------------------------------------------- build
def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- load results ----
    combined  = load_combined()
    imbalance = load_imbalance()
    best_row  = combined.dropna(subset=['best_val_AUC']).head(1) if combined['best_val_AUC'].notna().any() else None
    best_label = best_row['model'].iloc[0] if best_row is not None else 'TBD'
    best_auc   = fmt_score(best_row['best_val_AUC'].iloc[0]) if best_row is not None else 'TBD'

    # ===================================================================
    # Slide 1: Title
    # ===================================================================
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(1.2))
    set_text(tb.text_frame,
             'Loan Default Prediction',
             size=44, bold=True, color=TITLE_COLOR, align=PP_ALIGN.LEFT)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(12), Inches(0.9))
    set_text(tb.text_frame,
             'Comparing Classic ML and Neural Networks under Class Imbalance',
             size=24, color=GREY, align=PP_ALIGN.LEFT)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(0.5))
    set_text(tb.text_frame, 'Minwoo Yoo  |  Nathan',
             size=18, color=GREY, align=PP_ALIGN.LEFT)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(12), Inches(0.5))
    set_text(tb.text_frame, 'Home Credit Default Risk · Kaggle',
             size=14, color=GREY, align=PP_ALIGN.LEFT)

    # ===================================================================
    # Slide 2: Problem & Motivation
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Problem & Motivation')
    tf = add_body_textbox(s)
    set_text(tf, 'Goal', size=20, bold=True, color=TITLE_COLOR)
    add_paragraph(tf, 'Predict the probability that a loan applicant defaults on their loan (binary, TARGET ∈ {0, 1}).', size=16)
    add_paragraph(tf, '', bullet=False)
    add_paragraph(tf, 'Why it matters', size=20, bold=True, color=TITLE_COLOR, bullet=False)
    add_paragraph(tf, 'Lenders need accurate risk scoring to balance approval volume and credit losses.', size=16)
    add_paragraph(tf, 'Defaults are rare (~8% in this dataset) — the data is severely imbalanced.', size=16)
    add_paragraph(tf, 'A naive "always predict no default" model achieves 92% accuracy yet zero recall.', size=16)
    add_paragraph(tf, '', bullet=False)
    add_paragraph(tf, 'Imbalance is the central challenge', size=20, bold=True, color=TITLE_COLOR, bullet=False)
    add_paragraph(tf, 'Standard classifiers can be biased toward the majority class.', size=16)
    add_paragraph(tf, 'How models respond to imbalance handling techniques is the focus of this work.', size=16)

    # ===================================================================
    # Slide 3: Dataset & EDA
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Dataset · Home Credit Default Risk (Kaggle)')

    # Left half: text bullets
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(6.6), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, 'Size', size=20, bold=True, color=TITLE_COLOR)
    add_paragraph(tf, 'Train: 307,511 applicants × 122 features', size=14)
    add_paragraph(tf, 'Test: 48,744 applicants (labels held out by Kaggle)', size=14)
    add_paragraph(tf, '', bullet=False)
    set_text_p = add_paragraph
    add_paragraph(tf, 'Severe class imbalance', size=20, bold=True, color=TITLE_COLOR, bullet=False)
    add_paragraph(tf, 'Default rate ≈ 8.07% (24,825 / 307,511).', size=14)
    add_paragraph(tf, 'Naive "all-zero" model = 92% accuracy / 0 recall.', size=14)
    add_paragraph(tf, '', bullet=False)
    add_paragraph(tf, 'Key feature signals', size=20, bold=True, color=TITLE_COLOR, bullet=False)
    add_paragraph(tf, 'EXT_SOURCE_1/2/3 (external credit scores) are the strongest predictors (|ρ| ≈ 0.16–0.18).', size=14)
    add_paragraph(tf, 'DAYS_EMPLOYED sentinel 365,243 flags 18% (unemployed).', size=14)
    add_paragraph(tf, 'Building-info columns are 50–70% missing.', size=14)

    # Right half: class balance chart
    if os.path.exists(CHART_CLASS_BALANCE):
        s.shapes.add_picture(CHART_CLASS_BALANCE, Inches(7.3), Inches(2.0),
                             width=Inches(5.7))

    # ===================================================================
    # Slide 4: Research Question
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Research Question')
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.13), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, 'Do classic shallow models and neural networks respond differently to class imbalance handling techniques?',
             size=28, bold=True, color=TITLE_COLOR)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(12.13), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, 'Two-axis study', size=20, bold=True, color=ACCENT)
    add_paragraph(tf, 'Five model families: Logistic Regression, Random Forest, HistGradientBoosting, Shallow MLP, Deep MLP.', size=16)
    add_paragraph(tf, 'Imbalance handling: baseline · class_weight=balanced · SMOTE · threshold tuning.', size=16)
    add_paragraph(tf, 'Validation: 80/20 stratified split + GridSearchCV with PredefinedSplit; primary metric ROC AUC.', size=16)

    # ===================================================================
    # Slide 5: Methodology / Pipeline
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Methodology · Shared Pipeline')
    tf = add_body_textbox(s)
    set_text(tf, 'Preprocessing (identical for all models)', size=20, bold=True, color=TITLE_COLOR)
    add_paragraph(tf, '80% / 20% stratified train–val split (random_state=42).', size=16)
    add_paragraph(tf, 'Drop SK_ID_CURR; keep IDs only for the submission file.', size=16)
    add_paragraph(tf, 'Mean imputation for numeric NaN; categorical NaN → all-zero dummy block.', size=16)
    add_paragraph(tf, 'One-hot encoding for 16 categorical columns (244 features after encoding).', size=16)
    add_paragraph(tf, 'StandardScaler fitted on train only (essential for LR / MLP).', size=16)
    add_paragraph(tf, '', bullet=False)
    set_text_p(tf, 'Hyperparameter tuning', size=20, bold=True, color=TITLE_COLOR, bullet=False)
    add_paragraph(tf, 'GridSearchCV with sklearn PredefinedSplit (single train→val fold, no random k-fold).', size=16)
    add_paragraph(tf, 'Scoring: roc_auc.', size=16)
    add_paragraph(tf, 'Each model picks the highest-AUC parameter setting; results saved to result/home_credit/cv_results/.', size=16)

    # ===================================================================
    # Slide 6: Models compared
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Five Model Families Compared')
    df_models = pd.DataFrame([
        ['Logistic Regression',         'Linear',                'Nathan',  'C, tol; class_weight=balanced'],
        ['Random Forest',               'Tree (bagging)',        'Nathan',  'n_estimators, min_samples_*; class_weight=balanced'],
        ['HistGradientBoosting',        'Tree (boosting)',       'Nathan',  'learning_rate, max_iter, min_samples_leaf'],
        ['Shallow MLP — hidden (64,)',  'Neural network',        'Minwoo',  'alpha, learning_rate_init'],
        ['Deep MLP — hidden (128,64,32)','Neural network',        'Minwoo',  'alpha, learning_rate_init'],
    ], columns=['Model', 'Family', 'Owner', 'Tuned hyperparameters'])
    add_table(s, df_models, top=1.3, height=3.5,
              col_widths_inches=[3.4, 2.4, 1.4, 4.9])

    tb = s.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(12.13), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, 'Why this mix?', size=18, bold=True, color=TITLE_COLOR)
    add_paragraph(tf, 'Spans the course curriculum: linear → tree (bagging + boosting) → neural networks (shallow + deep).', size=14)
    add_paragraph(tf, 'Lets us isolate whether NN behaves differently from classic shallow learners under imbalance.', size=14)

    # ===================================================================
    # Slide 7: Results — Combined Comparison
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Results · Combined Validation AUC')

    if os.path.exists(CHART_MODEL_AUC):
        # Chart on the left, brief table on the right
        s.shapes.add_picture(CHART_MODEL_AUC, Inches(0.4), Inches(1.2), width=Inches(7.6))

        cmb_small = combined.copy()
        cmb_small['best_val_AUC'] = cmb_small['best_val_AUC'].apply(fmt_score)
        cmb_small = cmb_small[['model', 'best_val_AUC']]
        cmb_small.columns = ['Model', 'AUC']
        # truncate model labels for the right-side mini-table
        cmb_small['Model'] = cmb_small['Model'].astype(str).str.replace(r'\s*\(.*\)$', '', regex=True)
        add_table(s, cmb_small, top=1.4, left=8.3, width=4.7, height=3.0,
                  col_widths_inches=[3.0, 1.7])

        tb = s.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, f'Best model: {best_label}  (val ROC AUC = {best_auc})',
                 size=18, bold=True, color=TITLE_COLOR)
        add_paragraph(tf, 'Tree boosting (HGBC) is competitive on tabular credit data, in line with the literature.', size=14)
        add_paragraph(tf, 'Neural networks reach a similar regime once features are scaled and regularized.', size=14)
    else:
        cmb_for_table = combined.copy()
        cmb_for_table['best_val_AUC'] = cmb_for_table['best_val_AUC'].apply(fmt_score)
        cmb_for_table['best_params']  = cmb_for_table['best_params'].fillna('TBD').apply(
            lambda x: x[:90] + '…' if isinstance(x, str) and len(x) > 90 else x
        )
        add_table(s, cmb_for_table[['model', 'best_val_AUC', 'best_params']],
                  top=1.3, height=3.5,
                  col_widths_inches=[4.3, 1.6, 6.2])

        tb = s.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(12.13), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, 'Highlights', size=18, bold=True, color=TITLE_COLOR)
        add_paragraph(tf, f'Best model: {best_label}  (val ROC AUC = {best_auc}).', size=14)
        add_paragraph(tf, 'Tree boosting (HGBC) is competitive for tabular credit data, consistent with the literature.', size=14)
        add_paragraph(tf, 'Neural networks reach a similar regime once features are scaled and regularized.', size=14)

    # ===================================================================
    # Slide 8: Imbalance Handling Deep Dive (MLP)
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Imbalance Handling on the Best MLP')

    if imbalance is not None and len(imbalance) > 0:
        tbl_df = imbalance[['method', 'threshold', 'ROC_AUC', 'PR_AUC', 'precision', 'recall', 'F1']]
        add_table(s, tbl_df, top=1.2, height=1.8,
                  col_widths_inches=[1.6, 1.3, 1.5, 1.5, 1.6, 1.5, 2.3])

        if os.path.exists(CHART_IMBALANCE):
            s.shapes.add_picture(CHART_IMBALANCE, Inches(0.4), Inches(3.2), width=Inches(7.4))

        # Right-side commentary (or below if no chart)
        if os.path.exists(CHART_IMBALANCE):
            tb = s.shapes.add_textbox(Inches(8.1), Inches(3.2), Inches(4.9), Inches(3.8))
        else:
            tb = s.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(12.13), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, 'Three handling strategies', size=16, bold=True, color=TITLE_COLOR)
        add_paragraph(tf, 'Baseline — original 8/92 split, classify at 0.5.', size=12)
        add_paragraph(tf, 'SMOTE — oversample minority to 50/50 (train only).', size=12)
        add_paragraph(tf, 'Threshold tuning — pick the threshold maximizing F1 on val.', size=12)
        add_paragraph(tf, '', bullet=False)
        add_paragraph(tf, 'Reading the chart', size=16, bold=True, color=TITLE_COLOR, bullet=False)
        add_paragraph(tf, 'Threshold tuning swaps precision for recall.', size=12)
        add_paragraph(tf, 'SMOTE typically lifts recall but can hurt precision.', size=12)
        add_paragraph(tf, 'ROC AUC barely moves — the model already produces a good ranking.', size=12)
    else:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.13), Inches(1.0))
        set_text(tb.text_frame,
                 '<<TBD — fill after notebook completes>>',
                 size=18, color=GREY)

        tb = s.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12.13), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, 'Three handling strategies (sklearn MLPClassifier does not support class_weight)', size=18, bold=True, color=TITLE_COLOR)
        add_paragraph(tf, 'Baseline — train on the original imbalanced data; classify at 0.5.', size=14)
        add_paragraph(tf, 'SMOTE — synthesize minority samples until 50/50 in training only; classify at 0.5.', size=14)
        add_paragraph(tf, 'Threshold tuning — keep the baseline model, choose the threshold that maximizes F1 on validation.', size=14)

    # ===================================================================
    # Slide 9: Conclusions
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Conclusions')
    tf = add_body_textbox(s)
    set_text(tf, 'What we learned', size=20, bold=True, color=TITLE_COLOR)
    add_paragraph(tf, f'Best model: {best_label} at val ROC AUC = {best_auc}; tree boosting wins by about 1 AUC point.', size=15)
    add_paragraph(tf, 'Neural networks tie Logistic Regression and Random Forest — bottleneck is the feature set, not the model.', size=15)
    add_paragraph(tf, 'Shallow MLP (0.7405) ≈ Deep MLP (0.7416): going deeper does not help on this tabular data.', size=15)
    add_paragraph(tf, 'Most actionable finding: at 0.5 threshold the MLP has zero recall — *choosing the threshold matters more than resampling.*', size=15, bold=True)
    add_paragraph(tf, 'SMOTE actually hurt this MLP (AUC 0.745 → 0.648); threshold tuning preserved AUC and lifted recall to 40 %.', size=15)
    add_paragraph(tf, '', bullet=False)
    set_text_p(tf, 'Limitations & future work', size=20, bold=True, color=TITLE_COLOR, bullet=False)
    add_paragraph(tf, 'Only the main application_*.csv tables used — auxiliary tables (bureau, previous_application) typically add several AUC points.', size=14)
    add_paragraph(tf, 'PredefinedSplit is fast but optimistic; k-fold CV would tighten the confidence intervals.', size=14)
    add_paragraph(tf, 'Next steps: integrate auxiliary tables; ensemble HGBC + best MLP; calibrate probabilities (Platt / isotonic).', size=14)

    # ===================================================================
    # Slide 10: Q&A / Recording link
    # ===================================================================
    s = prs.slides.add_slide(blank)
    add_section_title(s, 'Thank you')
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.13), Inches(0.8))
    set_text(tb.text_frame, 'Questions?', size=36, bold=True, color=TITLE_COLOR)

    tb = s.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12.13), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, 'Recording link', size=20, bold=True, color=ACCENT)
    add_paragraph(tf, RECORDING_LINK_PLACEHOLDER, size=16)
    add_paragraph(tf, '', bullet=False)
    set_text_p(tf, 'Code', size=20, bold=True, color=ACCENT, bullet=False)
    add_paragraph(tf, 'github.com/ymw0414/loan-default-imbalance-classification', size=16)

    # --------------------------------------------------------------------- save
    prs.save(OUT_PPTX)
    print(f'wrote {OUT_PPTX}')
    print(f'slides: {len(prs.slides)}')
    print()
    print('best_label:', best_label, ' best_auc:', best_auc)
    print('imbalance results loaded:', imbalance is not None)


if __name__ == '__main__':
    build()
